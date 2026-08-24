"""Dashboard component examples.

Dashboard utilities with streaming, interactive filtering, export capabilities,
and responsive layout examples.

Author: Portfolio Team
Date: 2024
"""

import io
import json
import logging
import webbrowser
from collections.abc import AsyncGenerator, Callable
from dataclasses import dataclass, field
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any

try:
    import redis.asyncio as aioredis
except Exception:
    aioredis = None
import dash
import dash_bootstrap_components as dbc
import numpy as np
import pandas as pd
import plotly.graph_objects as go
from dash import Input, Output, State, callback_context, dcc, html
from dash.exceptions import PreventUpdate

try:
    from dash_extensions import WebSocket
except Exception:  # pragma: no cover - optional dependency
    WebSocket = None
from flask_caching import Cache
from flask_cors import CORS

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:  # pragma: no cover - optional dependency
    Presentation = None
    Inches = None
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class DashboardConfig:
    """Configuration for dashboard"""

    # Application settings
    app_name: str = "Enhanced Dashboard"
    host: str = "127.0.0.1"
    port: int = 8050
    debug: bool = True

    # Appearance
    theme: str = "light"  # light, dark, auto
    primary_color: str = "#1f77b4"
    secondary_color: str = "#ff7f0e"
    font_family: str = "Roboto, sans-serif"

    # Features
    enable_realtime: bool = True
    enable_export: bool = True
    enable_filtering: bool = True
    enable_dark_mode: bool = True
    enable_mobile: bool = True
    enable_accessibility: bool = True

    # Performance
    cache_timeout: int = 300  # seconds
    data_refresh_interval: int = 5000  # milliseconds
    max_data_points: int = 10000
    enable_data_sampling: bool = True

    # Export settings
    export_formats: list[str] = field(
        default_factory=lambda: ["pdf", "pptx", "csv", "excel"]
    )
    pdf_page_size: str = "A4"
    pptx_template: str | None = None

    # Redis settings for real-time data
    redis_host: str = "localhost"
    redis_port: int = 6379
    redis_db: int = 0

    # WebSocket settings
    websocket_url: str = "ws://localhost:8050/ws"
    websocket_reconnect_interval: int = 5000


class EnhancedDashboard:
    """Enhanced Dashboard with real-time streaming, interactive filtering,
    export/export and responsive design baked in.

    The class wires up a Dash application (with Bootstrap styling), the
    complementary Flask server, caching, optional Redis-backed streaming, and
    component registries. Supply a custom :class:`DashboardConfig` to override
    host/port/theme details or feature toggles such as ``enable_realtime``.

    Examples:
    --------
    >>> dash_app = EnhancedDashboard()
    >>> dash_app.config.host, dash_app.config.port
    ('127.0.0.1', 8050)
    >>> isinstance(dash_app.app.layout, html.Div)
    True
    """

    def __init__(self, config: DashboardConfig | None = None):
        """Initialize the enhanced dashboard.

        Args:
            config: Dashboard configuration
        """
        self.config = config or DashboardConfig()

        # Initialize Dash app with Bootstrap theme
        self.app = dash.Dash(
            __name__,
            external_stylesheets=[
                dbc.themes.BOOTSTRAP,
                "https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css",
            ],
            suppress_callback_exceptions=True,
            meta_tags=[
                {"name": "viewport", "content": "width=device-width, initial-scale=1"}
            ],
        )

        # Configure Flask server
        self.server = self.app.server
        CORS(self.server)

        # Initialize cache
        self.cache = Cache(
            self.server,
            config={
                "CACHE_TYPE": "simple",
                "CACHE_DEFAULT_TIMEOUT": self.config.cache_timeout,
            },
        )
        # Bind memoized method after cache initialization to avoid NameError at import time.
        self.get_data = self.cache.memoize(timeout=self.config.cache_timeout)(
            self.get_data
        )

        # Initialize Redis for real-time data (optional)
        self.redis_client = None
        if self.config.enable_realtime:
            try:
                import redis

                self.redis_client = redis.Redis(
                    host=self.config.redis_host,
                    port=self.config.redis_port,
                    db=self.config.redis_db,
                    decode_responses=True,
                )
                self.redis_client.ping()
                logger.info("Connected to Redis for real-time data")
            except:
                logger.warning("Redis not available, real-time features disabled")
                self.redis_client = None

        # Component registry
        self.components = {}
        self.filters = {}
        self.data_sources = {}

        # Layout components
        self._build_layout()

        logger.info("Enhanced Dashboard initialized")

    def _build_layout(self):
        """Build the dashboard layout with responsive design."""
        # Navigation bar
        navbar = dbc.NavbarSimple(
            children=[
                dbc.NavItem(dbc.NavLink("Dashboard", href="/")),
                dbc.NavItem(dbc.NavLink("Analytics", href="/analytics")),
                dbc.NavItem(dbc.NavLink("Reports", href="/reports")),
                dbc.DropdownMenu(
                    children=[
                        dbc.DropdownMenuItem("Export PDF", id="export-pdf"),
                        dbc.DropdownMenuItem("Export PowerPoint", id="export-pptx"),
                        dbc.DropdownMenuItem("Export Excel", id="export-excel"),
                        dbc.DropdownMenuItem(divider=True),
                        dbc.DropdownMenuItem("Settings", id="settings"),
                    ],
                    nav=True,
                    in_navbar=True,
                    label="More",
                ),
            ],
            brand=self.config.app_name,
            brand_href="/",
            color="primary",
            dark=True,
            fluid=True,
            id="navbar",
        )

        # Dark mode toggle
        dark_mode_switch = (
            dbc.Switch(
                id="dark-mode-switch",
                label="Dark Mode",
                value=self.config.theme == "dark",
                className="mt-2",
            )
            if self.config.enable_dark_mode
            else html.Div()
        )

        # Filter panel
        filter_panel = (
            dbc.Card(
                [
                    dbc.CardHeader("Filters"),
                    dbc.CardBody(
                        [
                            dcc.DatePickerRange(
                                id="date-range-picker",
                                start_date=datetime.now() - timedelta(days=30),
                                end_date=datetime.now(),
                                display_format="YYYY-MM-DD",
                                className="mb-3",
                            ),
                            dcc.Dropdown(
                                id="category-filter",
                                options=[],
                                multi=True,
                                placeholder="Select categories...",
                                className="mb-3",
                            ),
                            dcc.Dropdown(
                                id="metric-filter",
                                options=[],
                                multi=True,
                                placeholder="Select metrics...",
                                className="mb-3",
                            ),
                            dbc.Button(
                                "Apply Filters",
                                id="apply-filters-btn",
                                color="primary",
                                className="w-100",
                            ),
                        ]
                    ),
                ],
                className="mb-3",
            )
            if self.config.enable_filtering
            else html.Div()
        )

        # Main content area
        content = dbc.Container(
            [
                dbc.Row(
                    [
                        dbc.Col([filter_panel], width=12, md=3, className="mb-3"),
                        dbc.Col(
                            [
                                dcc.Location(id="url", refresh=False),
                                html.Div(id="page-content"),
                                dcc.Interval(
                                    id="interval-component",
                                    interval=self.config.data_refresh_interval,
                                    n_intervals=0,
                                    disabled=not self.config.enable_realtime,
                                ),
                                WebSocket(id="ws", url=self.config.websocket_url)
                                if self.config.enable_realtime
                                else html.Div(),
                            ],
                            width=12,
                            md=9,
                        ),
                    ]
                ),
                dbc.Row([dbc.Col([dark_mode_switch], width=12, className="text-end")]),
            ],
            fluid=True,
            className="p-4",
        )

        # Accessibility features
        if self.config.enable_accessibility:
            content = html.Div(
                [
                    html.Div("Skip to main content", className="sr-only", tabIndex=0),
                    content,
                ],
                role="main",
                **{"aria-label": "Main Dashboard Content"},
            )

        # Set the app layout
        self.app.layout = html.Div(
            [
                dcc.Store(id="theme-store", storage_type="local"),
                dcc.Store(id="filter-store", storage_type="session"),
                dcc.Store(id="data-store", storage_type="memory"),
                navbar,
                content,
            ],
            id="main-container",
        )

    def register_component(self, component_id: str, component_func: Callable):
        """Register a visualization component.

        Args:
            component_id: Unique identifier for the component
            component_func: Function that returns a Dash component
        """
        self.components[component_id] = component_func
        logger.info(f"Registered component: {component_id}")

    def register_data_source(self, source_id: str, source_func: Callable):
        """Register a data source.

        Args:
            source_id: Unique identifier for the data source
            source_func: Function that returns data
        """
        self.data_sources[source_id] = source_func
        logger.info(f"Registered data source: {source_id}")

    def get_data(self, source_id: str, filters: dict | None = None) -> pd.DataFrame:
        """Get data from a registered source with caching.

        Args:
            source_id: Data source identifier
            filters: Optional filters to apply

        Returns:
            DataFrame with the requested data
        """
        if source_id not in self.data_sources:
            raise ValueError(f"Data source {source_id} not found")

        data = self.data_sources[source_id](filters)

        # Apply data sampling if needed
        if self.config.enable_data_sampling and len(data) > self.config.max_data_points:
            sample_size = min(self.config.max_data_points, len(data))
            data = data.sample(n=sample_size, random_state=42)
            logger.info(f"Sampled data from {len(data)} to {sample_size} points")

        return data

    def create_responsive_card(
        self, title: str, content: Any, card_id: str | None = None
    ) -> dbc.Card:
        """Create a responsive card component.

        Args:
            title: Card title
            content: Card content
            card_id: Optional card ID

        Returns:
            Responsive card component
        """
        return dbc.Card(
            [
                dbc.CardHeader(
                    html.H5(title, className="mb-0"),
                    className="d-flex justify-content-between align-items-center",
                ),
                dbc.CardBody(content),
            ],
            id=card_id,
            className="mb-3 shadow-sm",
        )

    def create_kpi_card(
        self,
        title: str,
        value: str | float,
        change: float | None = None,
        icon: str | None = None,
    ) -> dbc.Card:
        """Create a KPI card with optional change indicator.

        Args:
            title: KPI title
            value: KPI value
            change: Optional change percentage
            icon: Optional Font Awesome icon class

        Returns:
            KPI card component
        """
        change_indicator = html.Div()
        if change is not None:
            color = "success" if change >= 0 else "danger"
            arrow = "↑" if change >= 0 else "↓"
            change_indicator = html.Div(
                f"{arrow} {abs(change):.1f}%", className=f"text-{color} small"
            )

        icon_element = html.Div()
        if icon:
            icon_element = html.I(className=f"{icon} fa-2x text-primary mb-2")

        return dbc.Card(
            dbc.CardBody(
                [
                    icon_element,
                    html.H6(title, className="text-muted"),
                    html.H3(value, className="mb-0"),
                    change_indicator,
                ]
            ),
            className="text-center h-100 shadow-sm",
        )

    async def stream_data(self, channel: str) -> AsyncGenerator:
        """Stream real-time data from Redis.

        Args:
            channel: Redis channel to subscribe to

        Yields:
            Real-time data updates
        """
        if not self.redis_client:
            return

        async with aioredis.create_redis_pool(
            f"redis://{self.config.redis_host}:{self.config.redis_port}"
        ) as redis:
            channel_obj = await redis.subscribe(channel)
            async for message in channel_obj[0].iter():
                yield json.loads(message.decode("utf-8"))

    def export_to_pdf(
        self, figures: list[go.Figure], filename: str = "dashboard_report.pdf"
    ) -> str:
        """Export dashboard to PDF.

        Args:
            figures: List of Plotly figures to export
            filename: Output filename

        Returns:
            Path to the generated PDF
        """
        output_path = Path(filename)

        # Create PDF document
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18,
        )

        # Container for the 'Flowable' objects
        elements = []

        # Add title
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=24,
            textColor=colors.HexColor(self.config.primary_color),
            spaceAfter=30,
            alignment=1,  # Center
        )

        elements.append(Paragraph(f"{self.config.app_name} Report", title_style))
        elements.append(
            Paragraph(
                f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                styles["Normal"],
            )
        )
        elements.append(Spacer(1, 0.5 * inch))

        # Add figures
        for i, fig in enumerate(figures):
            # Convert Plotly figure to image
            img_bytes = fig.to_image(format="png", width=600, height=400)
            img = Image(io.BytesIO(img_bytes), width=6 * inch, height=4 * inch)

            elements.append(
                Paragraph(
                    f"Figure {i + 1}: {fig.layout.title.text}", styles["Heading2"]
                )
            )
            elements.append(img)
            elements.append(Spacer(1, 0.3 * inch))

            if (i + 1) % 2 == 0 and i < len(figures) - 1:
                elements.append(PageBreak())

        # Build PDF
        doc.build(elements)
        logger.info(f"PDF exported to {output_path}")

        return str(output_path)

    def export_to_powerpoint(
        self, figures: list[go.Figure], filename: str = "dashboard_presentation.pptx"
    ) -> str:
        """Export dashboard to PowerPoint.

        Args:
            figures: List of Plotly figures to export
            filename: Output filename

        Returns:
            Path to the generated PowerPoint
        """
        # Create presentation
        if self.config.pptx_template:
            prs = Presentation(self.config.pptx_template)
        else:
            prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"{self.config.app_name} Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Add figure slides
        for i, fig in enumerate(figures):
            # Add slide with title and content layout
            bullet_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Add title
            title = slide.shapes.title
            title.text = f"{fig.layout.title.text}"

            # Convert figure to image
            img_bytes = fig.to_image(format="png", width=800, height=600)
            img_stream = io.BytesIO(img_bytes)

            # Add image to slide
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            pic = slide.shapes.add_picture(img_stream, left, top, width, height)

        # Save presentation
        output_path = Path(filename)
        prs.save(str(output_path))
        logger.info(f"PowerPoint exported to {output_path}")

        return str(output_path)

    def setup_callbacks(self):
        """Setup dashboard callbacks for interactivity."""
        # Dark mode toggle
        if self.config.enable_dark_mode:

            @self.app.callback(
                [Output("main-container", "className"), Output("theme-store", "data")],
                [Input("dark-mode-switch", "value")],
                [State("theme-store", "data")],
            )
            def toggle_dark_mode(dark_mode, stored_theme):
                theme = "dark" if dark_mode else "light"
                container_class = "dark-theme" if dark_mode else "light-theme"
                return container_class, theme

        # Page routing
        @self.app.callback(
            Output("page-content", "children"), [Input("url", "pathname")]
        )
        def display_page(pathname):
            if pathname == "/analytics":
                return self._create_analytics_page()
            elif pathname == "/reports":
                return self._create_reports_page()
            else:
                return self._create_dashboard_page()

        # Real-time data updates
        if self.config.enable_realtime:

            @self.app.callback(
                Output("data-store", "data"),
                [Input("interval-component", "n_intervals"), Input("ws", "message")],
            )
            def update_data(n, ws_message):
                ctx = callback_context
                if not ctx.triggered:
                    raise PreventUpdate

                trigger_id = ctx.triggered[0]["prop_id"].split(".")[0]

                if trigger_id == "ws" and ws_message:
                    # Handle WebSocket message
                    return json.loads(ws_message["data"])
                else:
                    # Regular interval update
                    # Fetch latest data from data sources
                    data = {}
                    for source_id in self.data_sources:
                        try:
                            data[source_id] = self.get_data(source_id).to_dict()
                        except:
                            pass
                    return data

        # Export callbacks
        if self.config.enable_export:

            @self.app.callback(
                Output("export-status", "children"),
                [
                    Input("export-pdf", "n_clicks"),
                    Input("export-pptx", "n_clicks"),
                    Input("export-excel", "n_clicks"),
                ],
            )
            def handle_export(pdf_clicks, pptx_clicks, excel_clicks):
                ctx = callback_context
                if not ctx.triggered:
                    raise PreventUpdate

                button_id = ctx.triggered[0]["prop_id"].split(".")[0]

                # Get current figures from the dashboard
                figures = self._collect_current_figures()

                if button_id == "export-pdf":
                    filepath = self.export_to_pdf(figures)
                    return dbc.Alert(
                        f"PDF exported successfully to {filepath}",
                        color="success",
                        dismissable=True,
                    )
                elif button_id == "export-pptx":
                    filepath = self.export_to_powerpoint(figures)
                    return dbc.Alert(
                        f"PowerPoint exported successfully to {filepath}",
                        color="success",
                        dismissable=True,
                    )
                elif button_id == "export-excel":
                    # Export data to Excel
                    filepath = self._export_to_excel()
                    return dbc.Alert(
                        f"Excel exported successfully to {filepath}",
                        color="success",
                        dismissable=True,
                    )

                raise PreventUpdate

    def _create_dashboard_page(self) -> html.Div:
        """Create the main dashboard page."""
        return dbc.Container(
            [
                dbc.Row(
                    [
                        dbc.Col(
                            self.create_kpi_card(
                                "Total Revenue",
                                "$1.2M",
                                change=12.5,
                                icon="fas fa-dollar-sign",
                            ),
                            width=6,
                            md=3,
                            className="mb-3",
                        ),
                        dbc.Col(
                            self.create_kpi_card(
                                "Active Users", "8,543", change=5.2, icon="fas fa-users"
                            ),
                            width=6,
                            md=3,
                            className="mb-3",
                        ),
                        dbc.Col(
                            self.create_kpi_card(
                                "Conversion Rate",
                                "3.2%",
                                change=-1.1,
                                icon="fas fa-chart-line",
                            ),
                            width=6,
                            md=3,
                            className="mb-3",
                        ),
                        dbc.Col(
                            self.create_kpi_card(
                                "Avg. Order Value",
                                "$142",
                                change=8.7,
                                icon="fas fa-shopping-cart",
                            ),
                            width=6,
                            md=3,
                            className="mb-3",
                        ),
                    ]
                ),
                dbc.Row(
                    [
                        dbc.Col(
                            html.Div(id="main-chart-container"),
                            width=12,
                            lg=8,
                            className="mb-3",
                        ),
                        dbc.Col(
                            html.Div(id="side-charts-container"),
                            width=12,
                            lg=4,
                            className="mb-3",
                        ),
                    ]
                ),
                html.Div(id="export-status"),
            ],
            fluid=True,
        )

    def _create_analytics_page(self) -> html.Div:
        """Create the analytics page."""
        return dbc.Container(
            [
                html.H2("Analytics Dashboard"),
                html.Hr(),
                dbc.Row([dbc.Col(html.Div(id="analytics-content"), width=12)]),
            ],
            fluid=True,
        )

    def _create_reports_page(self) -> html.Div:
        """Create the reports page."""
        return dbc.Container(
            [
                html.H2("Reports"),
                html.Hr(),
                dbc.Row(
                    [
                        dbc.Col(
                            [
                                dbc.Card(
                                    [
                                        dbc.CardHeader("Generate Report"),
                                        dbc.CardBody(
                                            [
                                                dbc.Label("Report Type"),
                                                dcc.Dropdown(
                                                    id="report-type",
                                                    options=[
                                                        {
                                                            "label": "Executive Summary",
                                                            "value": "executive",
                                                        },
                                                        {
                                                            "label": "Detailed Analysis",
                                                            "value": "detailed",
                                                        },
                                                        {
                                                            "label": "Custom Report",
                                                            "value": "custom",
                                                        },
                                                    ],
                                                    value="executive",
                                                ),
                                                html.Br(),
                                                dbc.Label("Date Range"),
                                                dcc.DatePickerRange(
                                                    id="report-date-range",
                                                    start_date=datetime.now()
                                                    - timedelta(days=30),
                                                    end_date=datetime.now(),
                                                ),
                                                html.Br(),
                                                html.Br(),
                                                dbc.Button(
                                                    "Generate Report",
                                                    id="generate-report-btn",
                                                    color="primary",
                                                    className="w-100",
                                                ),
                                            ]
                                        ),
                                    ]
                                )
                            ],
                            width=12,
                            md=4,
                        ),
                        dbc.Col(html.Div(id="report-preview"), width=12, md=8),
                    ]
                ),
            ],
            fluid=True,
        )

    def _collect_current_figures(self) -> list[go.Figure]:
        """Collect all current figures from the dashboard."""
        # This is a placeholder - in a real implementation,
        # you would collect actual figures from the dashboard
        figures = []

        # Create sample figures
        fig1 = go.Figure(
            data=[go.Scatter(x=[1, 2, 3, 4], y=[10, 11, 12, 13])],
            layout=go.Layout(title="Sample Chart 1"),
        )
        figures.append(fig1)

        fig2 = go.Figure(
            data=[go.Bar(x=["A", "B", "C"], y=[1, 3, 2])],
            layout=go.Layout(title="Sample Chart 2"),
        )
        figures.append(fig2)

        return figures

    def _export_to_excel(self, filename: str = "dashboard_data.xlsx") -> str:
        """Export data to Excel."""
        output_path = Path(filename)

        with pd.ExcelWriter(output_path, engine="xlsxwriter") as writer:
            for source_id in self.data_sources:
                try:
                    df = self.get_data(source_id)
                    df.to_excel(writer, sheet_name=source_id[:31], index=False)
                except:
                    pass

        logger.info(f"Excel exported to {output_path}")
        return str(output_path)

    def run(self):
        """Run the dashboard application."""
        self.setup_callbacks()

        logger.info(
            f"Starting dashboard on http://{self.config.host}:{self.config.port}"
        )

        # Open browser if in debug mode
        if self.config.debug:
            webbrowser.open(f"http://{self.config.host}:{self.config.port}")

        self.app.run_server(
            host=self.config.host, port=self.config.port, debug=self.config.debug
        )


# Utility functions for common visualizations
def create_time_series_chart(
    df: pd.DataFrame,
    x_col: str,
    y_cols: list[str],
    title: str = "Time Series",
    animation: bool = False,
) -> go.Figure:
    """Create an interactive time series chart.

    Args:
        df: DataFrame with data
        x_col: Column name for x-axis
        y_cols: List of column names for y-axis
        title: Chart title
        animation: Enable animation

    Returns:
        Plotly figure
    """
    fig = go.Figure()

    for col in y_cols:
        fig.add_trace(
            go.Scatter(
                x=df[x_col],
                y=df[col],
                mode="lines+markers",
                name=col,
                hovertemplate="%{x}<br>%{y:.2f}<extra></extra>",
            )
        )

    fig.update_layout(
        title=title,
        xaxis_title=x_col,
        yaxis_title="Value",
        hovermode="x unified",
        showlegend=True,
        template="plotly_white",
    )

    if animation:
        fig.update_layout(
            updatemenus=[
                dict(
                    type="buttons",
                    showactive=False,
                    buttons=[
                        dict(
                            label="Play",
                            method="animate",
                            args=[None, {"frame": {"duration": 100}}],
                        )
                    ],
                )
            ]
        )

    return fig


def create_heatmap(df: pd.DataFrame, title: str = "Heatmap") -> go.Figure:
    """Create an interactive heatmap.

    Args:
        df: DataFrame with data
        title: Chart title

    Returns:
        Plotly figure
    """
    fig = go.Figure(
        data=go.Heatmap(
            z=df.values,
            x=df.columns,
            y=df.index,
            colorscale="Viridis",
            hovertemplate="%{x}<br>%{y}<br>Value: %{z:.2f}<extra></extra>",
        )
    )

    fig.update_layout(title=title, template="plotly_white")

    return fig


if __name__ == "__main__":
    # Example usage
    config = DashboardConfig(
        app_name="Data Science Portfolio Dashboard",
        enable_realtime=True,
        enable_dark_mode=True,
        enable_export=True,
    )

    dashboard = EnhancedDashboard(config)

    # Register sample data source
    def sample_data_source(filters=None):
        """Generate a synthetic daily revenue dataset for demo purposes."""
        dates = pd.date_range(start="2024-01-01", end="2024-01-31", freq="D")
        data = pd.DataFrame(
            {
                "date": dates,
                "revenue": np.random.uniform(10000, 50000, len(dates)),
                "users": np.random.uniform(1000, 5000, len(dates)),
                "conversion": np.random.uniform(0.01, 0.05, len(dates)),
            }
        )
        return data

    dashboard.register_data_source("sample_data", sample_data_source)

    # Run dashboard
    dashboard.run()
