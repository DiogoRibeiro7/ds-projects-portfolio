"""Enhanced Web Dashboard Application with Real-time Streaming and Interactive Features.
Includes REST API, WebSocket support, and comprehensive dashboard capabilities.
"""

import base64
import io
import json
import logging
import threading
import time
from dataclasses import asdict, dataclass
from datetime import datetime, timedelta

import numpy as np
import pandas as pd
import plotly.graph_objs as go
import plotly.io as pio
import redis
from flask import Flask, jsonify, render_template, request, send_file, session
from flask_cors import CORS
from flask_jwt_extended import (
    JWTManager,
    create_access_token,
    jwt_required,
)
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_socketio import SocketIO, emit

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)
app.config["SECRET_KEY"] = "your-secret-key-here"
app.config["JWT_SECRET_KEY"] = "jwt-secret-key-here"
app.config["JWT_ACCESS_TOKEN_EXPIRES"] = timedelta(hours=1)

# Initialize extensions
socketio = SocketIO(app, cors_allowed_origins="*", async_mode="threading")
CORS(app)
jwt = JWTManager(app)

# Rate limiting
limiter = Limiter(
    app=app, key_func=get_remote_address, default_limits=["200 per day", "50 per hour"]
)

# Redis connection for caching (optional)
try:
    redis_client = redis.Redis(host="localhost", port=6379, db=0, decode_responses=True)
    redis_client.ping()
    REDIS_AVAILABLE = True
except:
    redis_client = None
    REDIS_AVAILABLE = False
    logger.warning("Redis not available, using in-memory cache")

# In-memory cache fallback
memory_cache = {}


@dataclass
class DashboardConfig:
    """User-facing configuration toggles for the Flask dashboard.

    Attributes:
    ----------
    theme : str
        Visual theme (`"light"` or `"dark"`). Drives CSS selection.
    auto_refresh : bool
        When ``True`` the client requests updated widgets on a timer.
    refresh_interval : int
        Polling interval in milliseconds for live widgets (default 5000 ms).
    export_formats : List[str]
        Allowed export targets. Defaults to ``['pdf', 'png', 'pptx']``.
    mobile_breakpoint : int
        CSS breakpoint in pixels for switching to the mobile layout.

    Examples:
    --------
    >>> cfg = DashboardConfig(refresh_interval=10000, theme="dark")
    >>> cfg.export_formats
    ['pdf', 'png', 'pptx']
    """

    theme: str = "light"
    auto_refresh: bool = True
    refresh_interval: int = 5000
    export_formats: list[str] = None
    mobile_breakpoint: int = 768

    def __post_init__(self):
        """Populate default export formats when none are provided."""
        if self.export_formats is None:
            self.export_formats = ["pdf", "png", "pptx"]


class DataStreamer:
    """Handle real-time data streaming."""

    def __init__(self):
        self.streaming = False
        self.stream_thread = None
        self.data_generators = {}

    def start_stream(self, stream_id: str, interval: int = 1):
        """Start data streaming for a specific stream."""
        if stream_id not in self.data_generators:
            self.data_generators[stream_id] = self._create_data_generator(stream_id)

        def stream_data():
            while self.streaming and stream_id in self.data_generators:
                data = next(self.data_generators[stream_id])
                socketio.emit(
                    "data_update",
                    {
                        "stream_id": stream_id,
                        "data": data,
                        "timestamp": datetime.now().isoformat(),
                    },
                )
                time.sleep(interval)

        self.streaming = True
        self.stream_thread = threading.Thread(target=stream_data)
        self.stream_thread.daemon = True
        self.stream_thread.start()

    def stop_stream(self, stream_id: str):
        """Stop data streaming."""
        if stream_id in self.data_generators:
            del self.data_generators[stream_id]

    def _create_data_generator(self, stream_id: str):
        """Create a data generator for streaming."""
        while True:
            # Simulate real-time data
            data = {
                "value": np.random.randn() * 10 + 50,
                "timestamp": datetime.now().isoformat(),
                "category": np.random.choice(["A", "B", "C"]),
                "stream_id": stream_id,
            }
            yield data


class VisualizationEngine:
    """Create interactive visualizations."""

    @staticmethod
    def create_time_series(
        data: pd.DataFrame, title: str = "Time Series", animated: bool = False
    ) -> str:
        """Create interactive time series visualization."""
        fig = go.Figure()

        for column in data.columns:
            if column != "timestamp":
                fig.add_trace(
                    go.Scatter(
                        x=data.index if data.index.name else data.get("timestamp"),
                        y=data[column],
                        mode="lines+markers",
                        name=column,
                        hovertemplate="%{x}<br>%{y:.2f}<extra></extra>",
                    )
                )

        fig.update_layout(
            title=title,
            xaxis_title="Time",
            yaxis_title="Value",
            hovermode="x unified",
            template="plotly_white",
            showlegend=True,
            updatemenus=(
                [
                    {
                        "buttons": [
                            (
                                {
                                    "label": "Play",
                                    "method": "animate",
                                    "args": [
                                        None,
                                        {
                                            "frame": {"duration": 500},
                                            "fromcurrent": True,
                                        },
                                    ],
                                }
                                if animated
                                else {}
                            )
                        ],
                        "direction": "left",
                        "showactive": False,
                        "type": "buttons",
                        "x": 0.1,
                        "y": 1.1,
                    }
                ]
                if animated
                else []
            ),
        )

        if animated:
            # Add animation frames
            frames = []
            for i in range(len(data)):
                frame_data = data.iloc[: i + 1]
                frame = go.Frame(
                    data=[
                        go.Scatter(
                            x=frame_data.index,
                            y=frame_data[col],
                            mode="lines+markers",
                            name=col,
                        )
                        for col in data.columns
                        if col != "timestamp"
                    ]
                )
                frames.append(frame)
            fig.frames = frames

        return pio.to_html(fig, include_plotlyjs="cdn", div_id="time-series-chart")

    @staticmethod
    def create_heatmap(data: pd.DataFrame, title: str = "Heatmap") -> str:
        """Create interactive heatmap."""
        fig = go.Figure(
            data=go.Heatmap(
                z=data.values,
                x=data.columns,
                y=data.index,
                colorscale="Viridis",
                hovertemplate="X: %{x}<br>Y: %{y}<br>Value: %{z}<extra></extra>",
            )
        )

        fig.update_layout(
            title=title,
            xaxis_title="Features",
            yaxis_title="Samples",
            template="plotly_white",
        )

        return pio.to_html(fig, include_plotlyjs="cdn", div_id="heatmap-chart")

    @staticmethod
    def create_3d_scatter(data: pd.DataFrame, title: str = "3D Scatter") -> str:
        """Create interactive 3D scatter plot."""
        if len(data.columns) < 3:
            raise ValueError("Need at least 3 columns for 3D scatter")

        fig = go.Figure(
            data=[
                go.Scatter3d(
                    x=data.iloc[:, 0],
                    y=data.iloc[:, 1],
                    z=data.iloc[:, 2],
                    mode="markers",
                    marker=dict(
                        size=5,
                        color=data.iloc[:, 2],
                        colorscale="Viridis",
                        showscale=True,
                    ),
                    hovertemplate="X: %{x:.2f}<br>Y: %{y:.2f}<br>Z: %{z:.2f}<extra></extra>",
                )
            ]
        )

        fig.update_layout(
            title=title,
            scene=dict(
                xaxis_title=data.columns[0],
                yaxis_title=data.columns[1],
                zaxis_title=data.columns[2],
            ),
            template="plotly_white",
        )

        return pio.to_html(fig, include_plotlyjs="cdn", div_id="scatter3d-chart")

    @staticmethod
    def create_dashboard_layout(charts: list[str], layout: str = "grid") -> str:
        """Create responsive dashboard layout."""
        if layout == "grid":
            html = '<div class="dashboard-grid">'
            for i, chart in enumerate(charts):
                html += (
                    f'<div class="chart-container" data-chart-id="{i}">{chart}</div>'
                )
            html += "</div>"
        elif layout == "stacked":
            html = '<div class="dashboard-stacked">'
            for i, chart in enumerate(charts):
                html += f'<div class="chart-container full-width" data-chart-id="{i}">{chart}</div>'
            html += "</div>"
        else:
            html = '<div class="dashboard-flex">'
            for i, chart in enumerate(charts):
                html += f'<div class="chart-container flex-item" data-chart-id="{i}">{chart}</div>'
            html += "</div>"

        return html


class ExportManager:
    """Handle export functionality."""

    @staticmethod
    def export_to_pdf(html_content: str, filename: str = "dashboard.pdf") -> bytes:
        """Export dashboard to PDF."""
        try:
            import pdfkit

            options = {
                "page-size": "A4",
                "margin-top": "0.75in",
                "margin-right": "0.75in",
                "margin-bottom": "0.75in",
                "margin-left": "0.75in",
                "enable-javascript": None,
                "javascript-delay": 1000,
                "no-stop-slow-scripts": None,
            }

            pdf = pdfkit.from_string(html_content, False, options=options)
            return pdf
        except ImportError:
            logger.error("pdfkit not installed")
            return None

    @staticmethod
    def export_to_pptx(charts: list[dict], filename: str = "dashboard.pptx") -> bytes:
        """Export dashboard to PowerPoint."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()

            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = "Dashboard Export"
            subtitle = title_slide.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

            # Add chart slides
            for chart_data in charts:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = chart_data.get("title", "Chart")

                # Add chart image if available
                if "image" in chart_data:
                    left = Inches(1)
                    top = Inches(2)
                    pic = slide.shapes.add_picture(
                        io.BytesIO(base64.b64decode(chart_data["image"])),
                        left,
                        top,
                        width=Inches(8),
                    )

            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.getvalue()

        except ImportError:
            logger.error("python-pptx not installed")
            return None


# Data streamer instance
data_streamer = DataStreamer()

# Visualization engine instance
viz_engine = VisualizationEngine()

# Export manager instance
export_manager = ExportManager()


# ============== ROUTES ==============


@app.route("/")
def index():
    """Main dashboard page."""
    return render_template("dashboard.html")


def create_app() -> Flask:
    """Factory for tests and external integrations."""
    return app


@app.route("/api/auth/login", methods=["POST"])
@limiter.limit("5 per minute")
def login():
    """API authentication endpoint."""
    username = request.json.get("username")
    password = request.json.get("password")

    # Simple authentication (replace with real authentication)
    if username and password == "password":
        access_token = create_access_token(identity=username)
        return jsonify(access_token=access_token), 200

    return jsonify(msg="Invalid credentials"), 401


@app.route("/api/data/<dataset_name>")
@jwt_required()
@limiter.limit("100 per minute")
def get_data(dataset_name: str):
    """Serve cached dashboard data sets via the authenticated REST API.

    Parameters
    ----------
    dataset_name : str
        Logical dataset key such as ``"churn_overview"``. A distinct Redis
        key (``data:<name>``) is created per dataset.

    Returns:
    -------
    flask.Response
        JSON payload containing the dataset plus cache headers. Responses are
        cached for 300 seconds (Redis when available, otherwise an in-process
        dictionary).

    Notes:
    -----
    This endpoint requires a valid JWT (``@jwt_required``) and is rate limited
    to 100 requests/minute. The fallback ``generate_sample_data`` helper is
    used when a live data source is not configured.

    Examples:
    --------
    >>> from flask.testing import FlaskClient
    >>> client = app.test_client()
    >>> token = create_access_token(identity="demo")
    >>> headers = {"Authorization": f"Bearer {token}"}
    >>> resp = client.get("/api/data/sales", headers=headers)
    >>> resp.status_code in (200, 401)  # depends on auth setup
    True
    """
    # Check cache first
    cache_key = f"data:{dataset_name}"

    if REDIS_AVAILABLE:
        cached = redis_client.get(cache_key)
        if cached:
            return json.loads(cached)
    elif cache_key in memory_cache:
        return memory_cache[cache_key]

    # Generate sample data (replace with real data source)
    data = generate_sample_data(dataset_name)

    # Cache the result
    result = jsonify(data)
    if REDIS_AVAILABLE:
        redis_client.setex(cache_key, 300, json.dumps(data))
    else:
        memory_cache[cache_key] = data

    return result


@app.route("/api/visualizations/<viz_type>", methods=["POST"])
@jwt_required()
def create_visualization(viz_type: str):
    """Create visualization based on type."""
    data = request.json.get("data")
    config = request.json.get("config", {})

    df = pd.DataFrame(data)

    if viz_type == "timeseries":
        html = viz_engine.create_time_series(
            df,
            title=config.get("title", "Time Series"),
            animated=config.get("animated", False),
        )
    elif viz_type == "heatmap":
        html = viz_engine.create_heatmap(df, title=config.get("title", "Heatmap"))
    elif viz_type == "3dscatter":
        html = viz_engine.create_3d_scatter(df, title=config.get("title", "3D Scatter"))
    else:
        return jsonify(error="Unknown visualization type"), 400

    return jsonify(html=html)


@app.route("/api/export/<format>", methods=["POST"])
@jwt_required()
def export_dashboard(format: str):
    """Export dashboard in specified format."""
    content = request.json.get("content")

    if format == "pdf":
        pdf_bytes = export_manager.export_to_pdf(content)
        if pdf_bytes:
            return send_file(
                io.BytesIO(pdf_bytes),
                mimetype="application/pdf",
                as_attachment=True,
                download_name="dashboard.pdf",
            )
    elif format == "pptx":
        charts = request.json.get("charts", [])
        pptx_bytes = export_manager.export_to_pptx(charts)
        if pptx_bytes:
            return send_file(
                io.BytesIO(pptx_bytes),
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                as_attachment=True,
                download_name="dashboard.pptx",
            )

    return jsonify(error="Export failed or format not supported"), 400


@app.route("/api/config", methods=["GET", "POST"])
@jwt_required()
def dashboard_config():
    """Get or update dashboard configuration."""
    if request.method == "GET":
        config = session.get("dashboard_config", asdict(DashboardConfig()))
        return jsonify(config)
    else:
        config_data = request.json
        session["dashboard_config"] = config_data
        return jsonify(success=True)


# ============== WebSocket Events ==============


@socketio.on("connect")
def handle_connect():
    """Handle client connection."""
    logger.info(f"Client connected: {request.sid}")
    emit("connected", {"data": "Connected to server"})


@socketio.on("disconnect")
def handle_disconnect():
    """Handle client disconnection."""
    logger.info(f"Client disconnected: {request.sid}")


@socketio.on("start_stream")
def handle_start_stream(data):
    """Start data streaming."""
    stream_id = data.get("stream_id", "default")
    interval = data.get("interval", 1)
    data_streamer.start_stream(stream_id, interval)
    emit("stream_started", {"stream_id": stream_id})


@socketio.on("stop_stream")
def handle_stop_stream(data):
    """Stop data streaming."""
    stream_id = data.get("stream_id", "default")
    data_streamer.stop_stream(stream_id)
    emit("stream_stopped", {"stream_id": stream_id})


@socketio.on("request_update")
def handle_request_update(data):
    """Handle update request from client."""
    dataset = data.get("dataset")
    update_data = generate_sample_data(dataset)
    emit("data_update", update_data)


# ============== Helper Functions ==============


def generate_sample_data(dataset_name: str) -> dict:
    """Generate sample data for demonstration."""
    np.random.seed(None)  # Random seed for variety

    if dataset_name == "timeseries":
        dates = pd.date_range(start="2024-01-01", periods=100, freq="H")
        data = {
            "timestamp": dates.strftime("%Y-%m-%d %H:%M:%S").tolist(),
            "value1": (np.cumsum(np.random.randn(100)) + 100).tolist(),
            "value2": (np.cumsum(np.random.randn(100)) + 50).tolist(),
            "value3": (np.sin(np.linspace(0, 4 * np.pi, 100)) * 20 + 80).tolist(),
        }
    elif dataset_name == "metrics":
        data = {
            "cpu_usage": float(np.random.uniform(20, 80)),
            "memory_usage": float(np.random.uniform(30, 70)),
            "disk_usage": float(np.random.uniform(40, 60)),
            "network_traffic": float(np.random.uniform(100, 1000)),
            "error_rate": float(np.random.uniform(0, 5)),
            "response_time": float(np.random.uniform(50, 500)),
        }
    elif dataset_name == "categories":
        categories = ["Product A", "Product B", "Product C", "Product D"]
        data = {
            "categories": categories,
            "sales": np.random.randint(100, 1000, len(categories)).tolist(),
            "revenue": np.random.uniform(10000, 100000, len(categories)).tolist(),
            "growth": np.random.uniform(-10, 30, len(categories)).tolist(),
        }
    else:
        # Default random data
        data = {
            "x": np.random.randn(50).tolist(),
            "y": np.random.randn(50).tolist(),
            "z": np.random.randn(50).tolist(),
            "category": np.random.choice(["A", "B", "C"], 50).tolist(),
        }

    return {
        "dataset": dataset_name,
        "data": data,
        "timestamp": datetime.now().isoformat(),
        "count": len(data.get("x", data.get("timestamp", []))),
    }


if __name__ == "__main__":
    socketio.run(app, debug=True, port=5000)
